"""Bouwt Som29.pptx: de uitlegvideo als PowerPoint met verschijn-op-klik-builds.

- Formules = PNG's uit genereer-assets.py (transparant)
- Cirkels, strepen, pijlen, chips = echte PowerPoint-vormen (zelf te bewerken)
- Animaties = timing-XML (python-pptx kan dat niet, dus injecteren we het)

Gebruik: python3 bouw-pptx.py  (vanuit video/powerpoint/)
"""
from PIL import Image
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt

matplotlib.rcParams["mathtext.fontset"] = "dejavuserif"

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor

GROEN = RGBColor.from_string("2D6A4F")
GROEN_LICHT = RGBColor.from_string("E8F5EE")
ROOD = RGBColor.from_string("C94F4A")
ROOD_LICHT = RGBColor.from_string("FBEEED")
TEKST = RGBColor.from_string("1A1A18")
GRIJS = RGBColor.from_string("7A7870")
RAND = RGBColor.from_string("E4E2D9")
WIT = RGBColor.from_string("FFFFFF")

SERIF = "DM Serif Display"
SANS = "DM Sans"

BREED = Inches(13.333)
HOOG = Inches(7.5)


def px_breedte(latex: str, pt: float) -> float:
    """Pixelbreedte van een mathtext-string (voor substring-posities)."""
    fig = plt.figure()
    t = fig.text(0, 0, latex, fontsize=pt)
    fig.canvas.draw()
    w = t.get_window_extent().width
    plt.close(fig)
    return w


def afbeelding(slide, pad: str, midden_x, top, breedte):
    """Plaats PNG met opgegeven breedte, horizontaal gecentreerd op midden_x."""
    im = Image.open(pad)
    ratio = im.height / im.width
    hoogte = Emu(int(breedte * ratio))
    return slide.shapes.add_picture(pad, Emu(int(midden_x - breedte / 2)), top, width=Emu(int(breedte))), hoogte


def tekstvak(slide, x, y, w, h, runs, pt, align=PP_ALIGN.CENTER, font=SANS, vet=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    for tekst, kleur in runs:
        r = p.add_run()
        r.text = tekst
        r.font.size = Pt(pt)
        r.font.name = font
        r.font.bold = vet
        r.font.color.rgb = kleur
    return tb


def chip(slide, x, y, w, h, tekst, kleur, vulling, pt=14):
    vorm = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    vorm.adjustments[0] = 0.5
    vorm.fill.solid()
    vorm.fill.fore_color.rgb = vulling
    vorm.line.fill.background()
    vorm.shadow.inherit = False
    tf = vorm.text_frame
    tf.margin_left = tf.margin_right = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = tekst
    r.font.size = Pt(pt)
    r.font.name = SANS
    r.font.bold = True
    r.font.color.rgb = kleur
    return vorm


def kaart(slide, x, y, w, h, vulling, randkleur=None):
    vorm = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    vorm.adjustments[0] = 0.12
    vorm.fill.solid()
    vorm.fill.fore_color.rgb = vulling
    vorm.shadow.inherit = False
    if randkleur is None:
        vorm.line.fill.background()
    else:
        vorm.line.color.rgb = randkleur
        vorm.line.width = Pt(1.75)
    return vorm


def ovaal(slide, x, y, w, h, kleur):
    vorm = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    vorm.fill.background()
    vorm.line.color.rgb = kleur
    vorm.line.width = Pt(2.5)
    vorm.shadow.inherit = False
    return vorm


def streep(slide, x1, y1, x2, y2, kleur):
    lijn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    lijn.line.color.rgb = kleur
    lijn.line.width = Pt(2.5)
    lijn.shadow.inherit = False
    return lijn


def pijl_omlaag(slide, midden_x, y):
    return tekstvak(slide, Emu(int(midden_x - Inches(0.3))), y, Inches(0.6), Inches(0.4), [("↓", GRIJS)], 22)


def logo(slide):
    tekstvak(
        slide,
        BREED - Inches(3.9),
        Inches(0.22),
        Inches(3.6),
        Inches(0.4),
        [("afgeleide", TEKST), ("oefenen", GROEN), (".nl", TEKST)],
        18,
        align=PP_ALIGN.RIGHT,
        font=SERIF,
    )


def aaf(slide, pose):
    hoogte = Inches(2.3)
    breedte = Emu(int(hoogte * 420 / 540))
    slide.shapes.add_picture(f"assets/aaf-{pose}.png", BREED - breedte - Inches(0.35), HOOG - hoogte - Inches(0.25), height=hoogte)


# ── Animatie: timing-XML met één (of meer) effecten per klik ──────────────
NS = 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'


def effect_xml(ctn_id: int, spid: int, node_type: str, richting: str) -> str:
    klasse = "entr" if richting == "in" else "exit"
    waarde = "visible" if richting == "in" else "hidden"
    return f"""<p:par {NS}><p:cTn id="{ctn_id}" presetID="1" presetClass="{klasse}" presetSubtype="0" fill="hold" grpId="0" nodeType="{node_type}"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="{ctn_id + 1}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="{waarde}"/></p:to></p:set></p:childTnLst></p:cTn></p:par>"""


def voeg_animaties_toe(slide, klikken):
    """klikken: lijst van kliks; elke klik is een lijst van (shape, 'in'|'out')."""
    ctn = 3
    klik_delen = []
    for klik in klikken:
        effecten = []
        for i, (vorm, richting) in enumerate(klik):
            node_type = "clickEffect" if i == 0 else "withEffect"
            effecten.append(effect_xml(ctn, vorm.shape_id, node_type, richting))
            ctn += 2
        binnen = f"""<p:par {NS}><p:cTn id="{ctn}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>{''.join(effecten)}</p:childTnLst></p:cTn></p:par>"""
        ctn += 1
        klik_delen.append(
            f"""<p:par {NS}><p:cTn id="{ctn}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>{binnen}</p:childTnLst></p:cTn></p:par>"""
        )
        ctn += 1

    blds = "".join(
        f'<p:bldP spid="{vorm.shape_id}" grpId="0"/>' for klik in klikken for (vorm, _r) in klik
    )
    xml = f"""<p:timing {NS}><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>{''.join(klik_delen)}</p:childTnLst></p:cTn><p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst><p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst><p:bldLst>{blds}</p:bldLst></p:timing>"""
    slide._element.append(parse_xml(xml))


# ── Presentatie ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = BREED
prs.slide_height = HOOG
blanco = prs.slide_layouts[6]

MIDden = Emu(int(Inches(6.0)))  # inhoud iets links van het midden (rechts staat Aaf)

# Substring-posities binnen de opgave (voor marker, cirkels en strepen).
# Let op: meet alleen prefixen die op een glyf eindigen (eindspaties vervallen);
# de startpositie van de "1" leiden we af via de losse breedte van "1".
OPGAVE = r"$m(q) = 1 - (3q^2 - 2)^2$"
PT_OPGAVE = 44
vol = px_breedte(OPGAVE, PT_OPGAVE)
r_q0 = px_breedte(r"$m($", PT_OPGAVE) / vol
r_q1 = px_breedte(r"$m(q$", PT_OPGAVE) / vol
r_e1 = px_breedte(r"$m(q) = 1$", PT_OPGAVE) / vol
r_e0 = r_e1 - 0.92 * px_breedte(r"$1$", PT_OPGAVE) / vol
r_min = px_breedte(r"$m(q) = 1 -$", PT_OPGAVE) / vol  # eind van het minteken
r_r0 = r_min + 0.012

OPGAVE_PROD = r"$m(q) = 1 - (3q^2-2) \cdot (3q^2-2)$"
PT_PROD = 40
vol_p = px_breedte(OPGAVE_PROD, PT_PROD)
p_e1 = px_breedte(r"$m(q) = 1$", PT_PROD) / vol_p
p_e0 = p_e1 - 0.92 * px_breedte(r"$1$", PT_PROD) / vol_p


def sub_x(img_x, img_w, ratio):
    return Emu(int(img_x + ratio * img_w))


# ── Gedeelde bouwstenen voor de slides ─────────────────────────────────────

def scene_titel(slide, runs, y=Inches(0.25), pt=32):
    return tekstvak(slide, Inches(1.0), y, Inches(10.0), Inches(0.7), runs, pt, font=SERIF, vet=True)


def stap_label(slide, y, tekst, x=Inches(0.75), w=Inches(2.9)):
    """Groen stap-pilletje, zoals de StapLabel in de video."""
    return chip(slide, x, y, w, Inches(0.42), tekst, WIT, GROEN, 13)


def regelkaart(slide, y, asset, breedte=Inches(7.6), hoogte=Inches(0.85)):
    """Groene reminder-kaart met de regel erin; geeft (kaart, formule) terug."""
    kaart_x = Emu(int(MIDden - breedte / 2))
    vlak = kaart(slide, kaart_x, y, breedte, hoogte, GROEN_LICHT)
    form, _h = afbeelding(slide, asset, MIDden, Emu(int(y + Inches(0.2))), Emu(int(breedte - Inches(1.0))))
    return vlak, form


def bijschrift(slide, y, tekst, x=Inches(3.4), w=Inches(6.4), align=PP_ALIGN.LEFT):
    return tekstvak(slide, x, y, w, Inches(0.34), [(tekst, GRIJS)], 12, align=align)


# ── Slide 1 · H2 · #29 ─────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blanco)
logo(s1)
aaf(s1, "wave")
titel29 = tekstvak(
    s1, Inches(3.0), Inches(2.1), Inches(6.0), Inches(1.1),
    [("H2 · ", TEKST), ("#29", GROEN)], 60, font=SERIF, vet=True,
)
opg1, _ = afbeelding(s1, "assets/opgave.png", MIDden, Inches(3.6), Inches(5.4))
voeg_animaties_toe(s1, [[(titel29, "in")], [(opg1, "in")]])

# ── Slide 2 · Stap 0: Analyseer de buitenste schil ─────────────────────────
s2 = prs.slides.add_slide(blanco)
logo(s2)
aaf(s2, "point")
titel2 = scene_titel(s2, [("Stap 0 · ", TEKST), ("Analyseer de buitenste schil", GROEN)], Inches(0.45))

eq_w = Inches(5.2)
eq_x = Emu(int(MIDden - eq_w / 2))
eq_y = Inches(1.75)

# marker achter de q (eerst toevoegen = achter de formule)
mark_x = sub_x(eq_x, eq_w, r_q0)
mark_w = Emu(int((r_q1 - r_q0) * eq_w + Inches(0.06)))
marker = kaart(s2, Emu(int(mark_x - Inches(0.03))), Emu(int(eq_y - Inches(0.04))), mark_w, Inches(0.62), ROOD_LICHT)

opg2, h2 = afbeelding(s2, "assets/opgave.png", MIDden, eq_y, eq_w)

c1_x = sub_x(eq_x, eq_w, r_e0)
c1_w = Emu(int((r_e1 - r_e0) * eq_w + Inches(0.20)))
cirkel1 = ovaal(s2, Emu(int(c1_x - Inches(0.05))), Emu(int(eq_y - Inches(0.16))), c1_w, Emu(int(h2 + Inches(0.32))), GROEN)
c2_x = sub_x(eq_x, eq_w, r_r0)
c2_w = Emu(int((1.0 - r_r0) * eq_w + Inches(0.22)))
cirkel2 = ovaal(s2, Emu(int(c2_x - Inches(0.03))), Emu(int(eq_y - Inches(0.16))), c2_w, Emu(int(h2 + Inches(0.32))), ROOD)

chip_som = chip(s2, Emu(int(MIDden - Inches(1.2))), Inches(3.35), Inches(2.4), Inches(0.55), "Somregel!", GROEN, GROEN_LICHT, 18)
kaart_som, form_som = regelkaart(s2, Inches(4.5), "assets/regel-som.png")

voeg_animaties_toe(
    s2,
    [
        [(opg2, "in")],
        [(marker, "in")],
        [(marker, "out")],
        [(cirkel1, "in")],
        [(cirkel2, "in")],
        [(chip_som, "in")],
        [(kaart_som, "in"), (form_som, "in")],
    ],
)

# ── Slide 3 · De somregel ──────────────────────────────────────────────────
s3 = prs.slides.add_slide(blanco)
logo(s3)
aaf(s3, "point")
scene_titel(s3, [("De ", TEKST), ("somregel", GROEN)])
regelkaart(s3, Inches(0.95), "assets/regel-som.png")

som_w = Inches(4.4)
som_x = Emu(int(MIDden - som_w / 2))
som_y = Inches(2.05)
opg3, h3 = afbeelding(s3, "assets/opgave.png", MIDden, som_y, som_w)
label_y = Emu(int(som_y + h3 + Inches(0.02)))
lab_g = tekstvak(
    s3, Emu(int(sub_x(som_x, som_w, r_e0) - Inches(0.5))), label_y, Inches(1.0), Inches(0.32),
    [("g(q)", GROEN)], 15, font=SERIF, vet=True,
)
lab_h = tekstvak(
    s3, Emu(int(sub_x(som_x, som_w, r_r0) + Inches(0.35))), label_y, Inches(1.0), Inches(0.32),
    [("h(q)", GROEN)], 15, font=SERIF, vet=True,
)

st1 = stap_label(s3, Inches(3.05), "Stap 1 · kies g en h")
g_def, _ = afbeelding(s3, "assets/som-g-def.png", Inches(4.6), Inches(3.0), Inches(1.35))
h_def, _ = afbeelding(s3, "assets/som-h-def.png", Inches(7.5), Inches(3.0), Inches(3.0))

st2 = stap_label(s3, Inches(3.75), "Stap 2 · bereken g′ en h′")
g_afg, _ = afbeelding(s3, "assets/som-g-afgeleide.png", Inches(4.7), Inches(3.7), Inches(1.55))
g_reden = bijschrift(s3, Inches(3.78), "(er zit geen q in)", Inches(5.7), Inches(3.0))

h_kwad, _ = afbeelding(s3, "assets/som-h-kwadraat.png", Inches(6.3), Inches(4.4), Inches(3.0))
pijl = pijl_omlaag(s3, Inches(6.3), Inches(4.95))
h_prod, _ = afbeelding(s3, "assets/som-h-product.png", Inches(6.3), Inches(5.45), Inches(4.6))
tip = bijschrift(s3, Inches(6.05), "handig: schrijf een kwadraat altijd op als de term keer zichzelf", Inches(3.4), Inches(5.8))

chip_prod = chip(s3, Emu(int(MIDden - Inches(1.4))), Inches(6.55), Inches(2.8), Inches(0.5), "Productregel!", ROOD, ROOD_LICHT, 16)

voeg_animaties_toe(
    s3,
    [
        [(st1, "in")],
        [(lab_g, "in"), (lab_h, "in")],
        [(g_def, "in"), (h_def, "in")],
        [(st2, "in")],
        [(g_afg, "in"), (g_reden, "in")],
        [(h_kwad, "in")],
        [(pijl, "in")],
        [(h_prod, "in")],
        [(tip, "in")],
        [(chip_prod, "in")],
    ],
)

# ── Slide 4 · De productregel ──────────────────────────────────────────────
s4 = prs.slides.add_slide(blanco)
logo(s4)
aaf(s4, "point")
scene_titel(s4, [("De ", TEKST), ("productregel", ROOD)])
regelkaart(s4, Inches(0.95), "assets/regel-product.png", Inches(9.2))

overname = bijschrift(s4, Inches(2.05), "de functie van de vorige pagina noemen we hier f", Inches(1.0), Inches(10.0), PP_ALIGN.CENTER)
f_def, _ = afbeelding(s4, "assets/prod-f-def.png", MIDden, Inches(2.4), Inches(5.2))

p_st1 = stap_label(s4, Inches(3.2), "Stap 1 · kies g en h")
p_g, _ = afbeelding(s4, "assets/g-def.png", Inches(5.1), Inches(3.15), Inches(2.1))
p_h, _ = afbeelding(s4, "assets/h-def.png", Inches(7.9), Inches(3.15), Inches(2.1))

p_st2 = stap_label(s4, Inches(3.85), "Stap 2 · bereken g′ en h′")
p_ga, _ = afbeelding(s4, "assets/g-afgeleide.png", Inches(5.1), Inches(3.8), Inches(1.6))
p_ha, _ = afbeelding(s4, "assets/h-afgeleide.png", Inches(7.6), Inches(3.8), Inches(1.6))

p_st3 = stap_label(s4, Inches(4.5), "Stap 3 · vul de formule in")
r1, _ = afbeelding(s4, "assets/stap3-symbolisch.png", Inches(7.4), Inches(4.45), Inches(5.0))
r2, _ = afbeelding(s4, "assets/stap3-ingevuld.png", Inches(7.4), Inches(5.05), Inches(5.6))
r3, _ = afbeelding(s4, "assets/stap3-samen.png", Inches(7.4), Inches(5.65), Inches(3.4))
r4, _ = afbeelding(s4, "assets/prod-antwoord.png", Inches(7.4), Inches(6.25), Inches(3.6))

voeg_animaties_toe(
    s4,
    [
        [(overname, "in")],
        [(f_def, "in")],
        [(p_st1, "in")],
        [(p_g, "in"), (p_h, "in")],
        [(p_st2, "in")],
        [(p_ga, "in"), (p_ha, "in")],
        [(p_st3, "in")],
        [(r1, "in")],
        [(r2, "in")],
        [(r3, "in")],
        [(r4, "in")],
    ],
)

# ── Slide 5 · Voeg alles samen ─────────────────────────────────────────────
s5 = prs.slides.add_slide(blanco)
logo(s5)
aaf(s5, "wave")
scene_titel(s5, [("Voeg alles ", TEKST), ("samen", GROEN)], Inches(0.5))
opg5, _ = afbeelding(s5, "assets/opgave.png", MIDden, Inches(1.45), Inches(4.4))
sam_h, _ = afbeelding(s5, "assets/samen-h-afgeleide.png", MIDden, Inches(2.6), Inches(4.0))
sam_sym, _ = afbeelding(s5, "assets/samen-symbolisch.png", MIDden, Inches(3.6), Inches(3.6))
sam_inv, _ = afbeelding(s5, "assets/samen-ingevuld.png", MIDden, Inches(4.6), Inches(4.4))
eind, _ = afbeelding(s5, "assets/eindantwoord.png", MIDden, Inches(5.7), Inches(5.2))

voeg_animaties_toe(
    s5,
    [
        [(opg5, "in")],
        [(sam_h, "in")],
        [(sam_sym, "in")],
        [(sam_inv, "in")],
        [(eind, "in")],
    ],
)

prs.save("Som29.pptx")
print("Som29.pptx geschreven —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
